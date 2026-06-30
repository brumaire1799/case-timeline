#!/usr/bin/env python3
"""
案件时间轴 PPT 生成器 v10（python-pptx）
- 支持多方当事人（parties 数组），每方独立颜色和位置
- 向后兼容 topParty/bottomParty 两方格式
- 圆角矩形文本框 + 真正箭头 + 统一黑体 + PowerPoint 兼容
用法: ~/miniconda3/bin/python3 generate_timeline_pptx.py events.json output.pptx
"""

import json, math, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.oxml.ns import qn

SLIDE_H = Inches(5.625)
MARGIN = Inches(0.30)
TIMELINE_Y = Inches(2.85)
DOT_R = Inches(0.045)
CONN_LEN = Cm(2.2)
CLEAR_TOP = Inches(0.07)
CLEAR_BOT = Inches(0.16)
CARD_H = Cm(2.5)

CARD_BOTTOM_Y = TIMELINE_Y - DOT_R - CLEAR_TOP - CONN_LEN  # 上方卡片下边沿（固定）
CARD_TOP_Y = TIMELINE_Y + DOT_R + CLEAR_BOT + CONN_LEN     # 下方卡片上边沿（固定）

COLOR_LINE = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT_NAME = '黑体'

# ── 多方 party 配置 ──
def hex_to_rgb(h):
    return RGBColor(int(h[1:3],16), int(h[3:5],16), int(h[5:7],16))

def build_parties(data):
    if 'parties' in data and len(data['parties']) > 0:
        return data['parties']
    tp = data.get('topParty', '原告')
    bp = data.get('bottomParty', '被告')
    return [
        {'name': tp, 'color': '#CC0000', 'position': 'top'},
        {'name': bp, 'color': '#333333', 'position': 'bottom'},
    ]

def parse_input(path):
    with open(path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    data['_parties'] = build_parties(data)
    data['_pcfg'] = {p['name']: {'color': p['color'], 'position': p['position']} for p in data['_parties']}
    bot = [p for p in data['_parties'] if p['position'] == 'bottom']
    data['_def'] = bot[-1] if bot else data['_parties'][-1]
    for e in data['events']:
        if 'party' not in e:
            e['party'] = data['_def']['name']
    return data

def resolve_party(data, party_str):
    """解析 party，支持"原告一、原告二"联合事件，返回 {colors, position}"""
    import re
    pc = data['_pcfg']
    if party_str in pc:
        return {'colors': [pc[party_str]['color']], 'position': pc[party_str]['position']}
    names = [s.strip() for s in re.split(r'[、,，]', party_str) if s.strip()]
    if len(names) > 1:
        cs = [pc[n]['color'] for n in names if n in pc]
        if cs:
            return {'colors': cs, 'position': pc.get(names[0], data['_def'])['position']}
    return {'colors': [data['_def']['color']], 'position': data['_def']['position']}

def parse_date(s):
    """解析日期字符串，返回天数（相对1970-01-01）"""
    import re
    m = re.match(r'(\d{4})年(\d{1,2})月(\d{1,2})日?', s)
    if m:
        from datetime import date
        return date(int(m[1]), int(m[2]), int(m[3])).toordinal()
    m = re.match(r'(\d{4})年(\d{1,2})月?', s)
    if m:
        from datetime import date
        return date(int(m[1]), int(m[2]), 15).toordinal()
    return 0

def calc_card_width(text):
    """计算在Pt(8)、最多6行内文本所需卡片宽度（字宽按9pt保守估算）"""
    cpl = max(1, math.ceil(len(text) / 6))
    text_w = cpl * int(9 * 12700)  # 9pt字宽（保守，实际渲染约8-8.5pt）
    return max(text_w + Inches(0.10), Inches(0.80))

def layout_positions(events):
    """计算每列的卡片宽度、中心x、卡片左边缘x，时间间距使用sqrt"""
    card_wids = [calc_card_width(e['event']) for e in events]
    xs = [MARGIN]
    for i in range(len(events) - 1):
        a = parse_date(events[i]['date'])
        b = parse_date(events[i+1]['date'])
        days = max(abs(b - a), 1)
        gap = int(math.sqrt(days) * 0.01 * 914400)  # sqrt(days)*0.01 英寸
        gap = max(gap, Inches(0.06))  # 最小间隙
        xs.append(xs[-1] + card_wids[i] + gap)
    cxs = [x + w / 2 for x, w in zip(xs, card_wids)]
    total_w = xs[-1] + card_wids[-1] + MARGIN
    return card_wids, cxs, total_w

def add_arrow(slide, cx, from_y, to_y, color):
    """画真正的箭头线：连接器 + 末端三角箭头（先创后设坐标）"""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, 0, 0, 0, 0
    )
    # 显式设置端点坐标
    connector.begin_x = cx
    connector.begin_y = from_y
    connector.end_x = cx
    connector.end_y = to_y

    connector.line.color.rgb = color
    connector.line.width = Pt(1.5)

    # tailEnd 在 <p:spPr>/<a:ln> 内部
    spPr = connector._element.find(qn('p:spPr'))
    if spPr is None:
        spPr = connector._element.makeelement(qn('p:spPr'), {})
        connector._element.insert(0, spPr)
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = spPr.makeelement(qn('a:ln'), {'w': '19050'})
        spPr.append(ln)
    for old in ln.findall(qn('a:tailEnd')):
        ln.remove(old)
    tailEnd = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    ln.append(tailEnd)

def build_slide(prs, events, data, slide_w):
    title = data['title']
    parties = data['_parties']
    card_wids, cxs, total_w = layout_positions(events)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题（跟随幻灯片宽度）
    title_w = min(slide_w - MARGIN * 2 - Inches(2.2), Inches(50))
    tb = slide.shapes.add_textbox(MARGIN, Inches(0.05), title_w, Inches(0.35))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.font.name = FONT_NAME; p.alignment = PP_ALIGN.LEFT

    # 图例
    topP = [p for p in parties if p['position'] == 'top']
    botP = [p for p in parties if p['position'] == 'bottom']
    lx = slide_w - MARGIN - Inches(2.4)
    tb2 = slide.shapes.add_textbox(lx, Inches(0.02), Inches(2.4), Inches(0.38))
    tf2 = tb2.text_frame; tf2.word_wrap = True
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    if topP:
        r = p2.add_run(); r.text = '上方：'; r.font.size = Pt(7); r.font.bold = True
        r.font.color.rgb = RGBColor(0x66, 0x66, 0x66); r.font.name = FONT_NAME
        for pi in topP:
            r = p2.add_run(); r.text = f"● {pi['name']}  "
            r.font.size = Pt(7); r.font.bold = True; r.font.color.rgb = hex_to_rgb(pi['color']); r.font.name = FONT_NAME
    if botP:
        r = p2.add_run(); r.text = '下方：'; r.font.size = Pt(7); r.font.bold = True
        r.font.color.rgb = RGBColor(0x66, 0x66, 0x66); r.font.name = FONT_NAME
        for pi in botP:
            r = p2.add_run(); r.text = f"● {pi['name']}  "
            r.font.size = Pt(7); r.font.bold = True; r.font.color.rgb = hex_to_rgb(pi['color']); r.font.name = FONT_NAME

    # 时间轴主线
    main_line = slide.shapes.add_connector(
        MSO_CONNECTOR_TYPE.STRAIGHT, 0, 0, 0, 0
    )
    main_line.begin_x = MARGIN
    main_line.begin_y = TIMELINE_Y
    main_line.end_x = slide_w - MARGIN - Inches(0.05)
    main_line.end_y = TIMELINE_Y
    main_line.line.color.rgb = COLOR_LINE
    main_line.line.width = Pt(4.5)
    spPr = main_line._element.find(qn('p:spPr'))
    if spPr is None:
        spPr = main_line._element.makeelement(qn('p:spPr'), {})
        main_line._element.insert(0, spPr)
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        ln = spPr.makeelement(qn('a:ln'), {'w': '57150'})
        spPr.append(ln)
    tailEnd = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'lg', 'len': 'lg'})
    ln.append(tailEnd)

    for i, evt in enumerate(events):
        rp = resolve_party(data, evt['party'])
        is_top = rp['position'] == 'top'
        colors = [hex_to_rgb(c) for c in rp['colors']]
        color = colors[0]
        cx = cxs[i]
        box_w = card_wids[i]
        text_x = cx - box_w / 2
        if is_top:
            card_y = CARD_BOTTOM_Y - CARD_H
            arr_from = CARD_BOTTOM_Y
        else:
            card_y = CARD_TOP_Y
            arr_from = CARD_TOP_Y
        bg = RGBColor(0xFF, 0xF0, 0xF0) if is_top else RGBColor(0xF5, 0xF5, 0xF5)

        # 圆角矩形（固定2.5cm高度，固定8pt字号）
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, text_x, card_y, box_w, CARD_H
        )
        shape.fill.solid(); shape.fill.fore_color.rgb = bg
        shape.line.color.rgb = color; shape.line.width = Pt(1.5)
        prstGeom = shape._element.find(qn('a:prstGeom'))
        if prstGeom is not None:
            prstGeom.set('prst', 'roundRect')
            avLst = prstGeom.find(qn('a:avLst'))
            if avLst is None:
                avLst = prstGeom.makeelement(qn('a:avLst'), {})
                prstGeom.append(avLst)
            gd = avLst.makeelement(qn('a:gd'), {'name': 'adj', 'fmla': 'val 7500'})
            avLst.append(gd)

        # 文字（固定8pt，6行封顶）
        tf = shape.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.04)
        tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
        p = tf.paragraphs[0]; p.text = evt['event']
        p.font.size = Pt(8); p.font.bold = True; p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.line_spacing = Pt(10)
        p.alignment = PP_ALIGN.LEFT

        # 连接箭头
        arr_to = TIMELINE_Y - DOT_R - CLEAR_TOP if is_top else TIMELINE_Y + DOT_R + CLEAR_BOT
        add_arrow(slide, cx, arr_from, arr_to, color)

        # 圆点
        if len(colors) > 1:
            DOT_R2 = Inches(0.07)
            n = len(colors)
            FULL = 216.0
            ROT_BASE = 90.0 if n == 2 else 150.0
            for k in range(n):
                pie = slide.shapes.add_shape(
                    MSO_SHAPE.PIE, cx - DOT_R2, TIMELINE_Y - DOT_R2, DOT_R2 * 2, DOT_R2 * 2
                )
                pie.fill.solid(); pie.fill.fore_color.rgb = colors[k]
                pie.line.color.rgb = WHITE; pie.line.width = Pt(1)
                pie.adjustments[0] = 0
                pie.adjustments[1] = FULL / n
                pie.rotation = ROT_BASE + k * (360.0 / n)
        else:
            dot = slide.shapes.add_shape(
                MSO_SHAPE.OVAL, cx - DOT_R, TIMELINE_Y - DOT_R, DOT_R * 2, DOT_R * 2
            )
            dot.fill.solid(); dot.fill.fore_color.rgb = color
            dot.line.color.rgb = WHITE; dot.line.width = Pt(1.5)

        # 日期
        dtb = slide.shapes.add_textbox(text_x, TIMELINE_Y + DOT_R + Inches(0.02), box_w, Inches(0.14))
        dtf = dtb.text_frame; dtf.word_wrap = False
        dp = dtf.paragraphs[0]; dp.text = evt['date']
        dp.font.size = Pt(5.5); dp.font.bold = True
        dp.font.color.rgb = RGBColor(0x88, 0x88, 0x88)
        dp.font.name = FONT_NAME; dp.alignment = PP_ALIGN.CENTER

def main():
    if len(sys.argv) < 3:
        print("用法: ~/miniconda3/bin/python3 generate_timeline_pptx.py <events.json> <output.pptx>")
        sys.exit(1)
    data = parse_input(sys.argv[1])
    events = data['events']
    card_wids, _, total_w = layout_positions(events)
    slide_w = min(total_w, Inches(56))  # PPT最大宽度56"
    prs = Presentation()
    prs.slide_width = slide_w; prs.slide_height = SLIDE_H
    build_slide(prs, events, data, slide_w)
    prs.save(sys.argv[2])
    n = len(events)
    print(json.dumps({"status":"ok","slides":1,"events":n,"slide_w_in":round(slide_w/914400,2),"output":sys.argv[2]}))

if __name__ == '__main__':
    main()
